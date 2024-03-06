from sklearn.calibration import LabelEncoder
import streamlit as st
import joblib
import re
import pandas as pd
from sklearn.feature_extraction.text import TfidfVectorizer
from streamlit_option_menu import option_menu
import pdfplumber
from docx import Document
from pptx import Presentation
import base64
import json
# Charger le modèle
#tfidf_model = joblib.load('models/tf_idf_model.joblib')
# Charger le modèle et le vecteur TF-IDF
logistic_regression_model = joblib.load('models/tf_idf_model.joblib')
tfidf_vectorizer = joblib.load('models/tfidf_vectorizer.joblib')  # Assurez-vous d'avoir sauvegardé le vecteur TF-IDF pendant l'entraînement
label_encoder = joblib.load('models/label_encoder.joblib')
# hide_st_style = """
#             <style>
#             #MainMenu {visibility: hidden;}
#             footer {visibility: hidden;}
#             header {visibility: hidden;}
#             </style>
#             """
# st.markdown(hide_st_style, unsafe_allow_html=False)
# Fonction pour effectuer la prédiction
def predict_category(text):
    # Prétraitement du texte (assurez-vous que c'est le même que celui utilisé pour l'entraînement)
    text = re.sub('[^a-zA-Z]', ' ', text)
    text = text.lower()

    # Utiliser le modèle TF-IDF pour la prédiction
    text_tfidf = tfidf_vectorizer.transform([text])
    prediction = logistic_regression_model.predict(text_tfidf)[0]
    predicted_category = label_encoder.inverse_transform([prediction])[0]

    return predicted_category
def convert_to_json(text):
    try:
        # Convertir le texte en structure JSON
        json_data = {"text_content": text}
        json_result = json.dumps(json_data, indent=2)
        return json_result
    except Exception as e:
        return f"Erreur lors de la conversion en JSON : {str(e)}"

# Logo (optionnel)
#st.image("le-wagon-logo.png", width=200)  # Remplacez "path/to/your/logo.png" par le chemin de votre logo

#st.sidebar.image("le-wagon-logo.png", width=90)
logo_style = """
    <style>
        .logo-text {
            font-size: 30px;
            font-weight: bold;
            color: #FFFFFF;  /* Couleur du texte */
            background-color: #A52A2A;  /* Couleur de fond */
            padding: 10px;
            border-radius: 10px;
            text-align: center;
        }
    </style>
"""

    # Afficher le style
st.markdown(logo_style, unsafe_allow_html=True)

    # Afficher le texte du logo
st.markdown('<div class="logo-text">Talent-sort-IA</div>', unsafe_allow_html=True)
with st.sidebar:
    #st.sidebar.image("le-wagon-logo.png", width=300)

    selected = option_menu(
        menu_title="Main menu",
        options=["Home","Prédiction Catégorie","Contacts","Recherche","À propos"],
        icons=["house","book","envelope","search"],
        menu_icon="cast",
        default_index=0,
    )




def load_image(path):
    with open(path, 'rb') as f:
        data = f.read()
    encoded = base64.b64encode(data).decode()
    return encoded

def image_tag(path):
    encoded = load_image(path)
    tag = f'<img src="data:image/png;base64,{encoded}">'
    return tag

def background_image_style(path):
    encoded = load_image(path)
    style = f'''
    <style>
    .stApp {{
        background-image: url("data:image/png;base64,{encoded}");
        background-size: cover;
    }}
    </style>
    '''
    return style

image_path = 'back.jpg'


#if st.checkbox('Show background image', True):
st.write(background_image_style(image_path), unsafe_allow_html=True)
#Contenu principal en fonction de l'option sélectionnée
if selected == "Home":
    st.title("Bienvenue dans notre application de catégorisation des CVs")
    st.write("""
    Notre application vous permet de catégoriser les CVs en fonction de leur contenu textuel.
    Vous pouvez soumettre le texte d'un CV, et notre modèle prédictif utilisant la méthode TF-IDF
    sera en mesure de prédire la catégorie correspondante. Que vous soyez un professionnel du recrutement,
    un étudiant ou simplement curieux de voir comment cela fonctionne, n'hésitez pas à explorer et à utiliser
    notre outil convivial.

    Instructions :
    1. Choisissez l'option "Text" dans le menu pour entrer le texte du CV.
    2. Cliquez sur le bouton "Prédire la catégorie" pour obtenir la prédiction.
    3. Explorez d'autres fonctionnalités disponibles dans le menu pour importer des fichiers ou effectuer des recherches spécifiques.

    N'hésitez pas à profiter de notre application et à découvrir les diverses fonctionnalités qu'elle offre !
    """)
    # Ajoutez le contenu de la page d'accueil ici


elif selected == "Prédiction Catégorie":
    # Interface utilisateur Streamlit
    st.title('Categorisation des CVs')


    # Texte d'introduction
    st.write("Bienvenue dans notre application de catégorisation des CVs.")
    st.write(" Sélectionnez le format du votre  CV pour prédire sa catégorie.")
    selected = option_menu(
        menu_title=None,
        options=["Csv","Pdf","Word","Img","pptx","text","JSON"],
        icons=["filetype-csv","file-pdf","file-word","file-image","file-ppt","file-earmark-text"],
        menu_icon="cast",
        default_index=0,
        orientation="horizontal",
    )



    if selected == "text":
        # Zone de saisie utilisateur
        user_input = st.text_area("Entrez le texte du CV ici:")
        # Bouton de prédiction
        if st.button("Prédire la catégorie"):
        # Vérifier si l'utilisateur a saisi quelque chose
            if user_input:
                # Obtenir la prédiction
                prediction = predict_category(user_input)

                st.success(f"La catégorie prédite est: {prediction}")
            else:
                st.warning("Veuillez saisir le texte du CV.")
            #st.title("Recherche dans les CVs")
            #search_input = st.text_input("Rechercher dans les CVs", "")
            # Ajoutez le contenu de la page de recherche ici

    if selected == "Pdf":
        st.set_option('deprecation.showfileUploaderEncoding', False)
        uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")

        if uploaded_file is not None:
            # Lire le fichier PDF
            with pdfplumber.open(uploaded_file) as pdf:
            # Extraire le texte de chaque page
                text = ""
                for page_num in range(len(pdf.pages)):
                    page = pdf.pages[page_num]
                    text += page.extract_text()

            # Afficher le texte extrait
            st.write(text)
        if st.button("Prédire la catégorie"):
        # Vérifier si l'utilisateur a saisi quelque chose
            if text:
                    # Obtenir la prédiction
                prediction = predict_category(text)

                st.success(f"La catégorie prédite est: {prediction}")

    if selected == "Csv":
        st.set_option('deprecation.showfileUploaderEncoding', False)
        uploaded_file = st.file_uploader("Choose a CSV file", type="csv")

        if uploaded_file is not None:
            data = pd.read_csv(uploaded_file)
            st.write(data)
        if st.button("Prédire la catégorie"):
        # Vérifier si l'utilisateur a saisi quelque chose
            if data:
                    # Obtenir la prédiction
                prediction = predict_category(data)

                st.success(f"La catégorie prédite est: {prediction}")
    if selected == "Word":
        st.set_option('deprecation.showfileUploaderEncoding', False)
        uploaded_file = st.file_uploader("Choose a Word file", type="docx")

        if uploaded_file is not None:
            # Lire le fichier Word
            doc = Document(uploaded_file)

            # Extraire le texte de chaque paragraphe
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"

            # Afficher le texte extrait
            st.write(text)
        if st.button("Prédire la catégorie"):
        # Vérifier si l'utilisateur a saisi quelque chose
            if text:
                    # Obtenir la prédiction
                prediction = predict_category(text)

                st.success(f"La catégorie prédite est: {prediction}")

    if selected == "pptx":
        st.set_option('deprecation.showfileUploaderEncoding', False)
        uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")

        if uploaded_file is not None:
            # Lire le fichier PowerPoint
            presentation = Presentation(uploaded_file)

            # Extraire le texte de chaque diapositive
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            # Afficher le texte extrait
            st.write(text)
        if st.button("Prédire la catégorie"):
        # Vérifier si l'utilisateur a saisi quelque chose
            if text:
                    # Obtenir la prédiction
                prediction = predict_category(text)

                st.success(f"La catégorie prédite est: {prediction}")

    if selected == "JSON":
        # Zone de saisie utilisateur
        user_input = st.text_area("Entrez le texte du CV ici:")
        # Bouton de prédiction
        if st.button(" Convertir en JSON "):
        # Vérifier si l'utilisateur a saisi quelque chose
            if user_input:
                # Obtenir la prédiction
                text_json = convert_to_json(user_input)
                st.text(text_json)

elif selected == "Recherche":
    st.title("Recherche dans les CVs")

    st.set_option('deprecation.showfileUploaderEncoding', False)
    uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")

    if uploaded_file is not None:
        # Lire le fichier PDF
        with pdfplumber.open(uploaded_file) as pdf:
            # Extraire le texte de chaque page
            text = ""
            for page_num in range(len(pdf.pages)):
                page = pdf.pages[page_num]
                text += page.extract_text()

        # Afficher le texte extrait
        st.write(text)

        # Zone de saisie pour la recherche
        search_input = st.text_input("Rechercher dans les CVs", "")

        # Bouton de recherche
        if st.button("Rechercher"):
            # Vérifier si l'utilisateur a saisi quelque chose
            if search_input:
                # Effectuer la recherche dans le texte
                match = re.search(search_input, text)

                # Afficher les résultats de la recherche
                if match:
                    st.success(f"Résultat trouvé : '{match.group()}'")
                else:
                    st.warning("Aucun résultat trouvé.")
            else:
                st.warning("Veuillez saisir quelque chose dans la zone de recherche.")



elif selected == "À propos":
    st.title("À propos")
    st.write("Description de l'application, informations sur les auteurs, etc.")
    # Ajoutez le contenu de la page À propos ici

st.markdown(
    """
    <style>
        .big-font {
            font-size: 24px !important;
        }
    </style>
    """,
    unsafe_allow_html=True
)

#st.write('<span class="big-font">Prédiction de Catégorie</span>', unsafe_allow_html=True)



# Ajouter d'autres éléments si nécessaire
